"""
CPT EOC PowerPoint Generator
Matches the Cherry Pick Talent EOC Canva template:
- Black/white palette with olive green accents
- Poppins-style bold headings
- CPT "cpt" wordmark top right on content slides
- Dark and light slide alternation matching the deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import re

# ── BRAND COLOURS ─────────────────────────────────────────────────────────────
BLACK   = RGBColor(0x11, 0x11, 0x11)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OLIVE   = RGBColor(0x4A, 0x5A, 0x0A)
OLIVE_L = RGBColor(0x8A, 0x9E, 0x20)
CREAM   = RGBColor(0xF5, 0xF4, 0xEF)
GREY    = RGBColor(0x88, 0x88, 0x88)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)

W = Inches(13.33)   # WIDE layout width
H = Inches(7.5)     # WIDE layout height

def fmt_num(n):
    if n is None or str(n).strip() in ('', '-', 'None'): return '—'
    try:
        n = float(str(n).replace(',','').replace('$','').replace('%',''))
        if n >= 1_000_000: return f"{n/1_000_000:.1f}M"
        if n >= 1_000:     return f"{n/1_000:.1f}K"
        return f"{int(n):,}"
    except: return str(n)

def fmt_pct(n):
    if n is None or str(n).strip() in ('', '-', 'None'): return '—'
    try: return f"{float(str(n).replace('%','')):.1f}%"
    except: return str(n)

def fmt_currency(n):
    if n is None or str(n).strip() in ('', '-', 'None'): return '—'
    try: return f"${float(str(n).replace(',','').replace('$','')):,.0f}"
    except: return str(n)

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             font_name="Arial", italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color or BLACK
    return txBox

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt as PPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_cpt_mark(slide):
    """CPT wordmark top right — appears on all content slides"""
    add_text(slide, "cpt", 12.1, 0.15, 1.0, 0.5,
             size=22, bold=True, color=BLACK, align=PP_ALIGN.RIGHT, font_name="Georgia")

def add_cpt_mark_white(slide):
    add_text(slide, "cpt", 12.1, 0.15, 1.0, 0.5,
             size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font_name="Georgia")

def stat_card(slide, x, y, w, h, value, label, dark=True):
    """Metric card matching Campaign Snapshot slide style"""
    bg = BLACK if dark else WHITE
    fg = WHITE if dark else BLACK
    border = WHITE if dark else BLACK
    add_rect(slide, x, y, w, h, bg, border, 1)
    add_text(slide, value, x+0.1, y+0.15, w-0.2, h*0.55,
             size=28, bold=True, color=fg, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(slide, label, x+0.1, y+h*0.6, w-0.2, h*0.35,
             size=11, bold=False, color=fg, align=PP_ALIGN.CENTER)


def generate_pptx(campaign_data, output_path):
    """
    campaign_data dict keys:
      campaign_title, campaign_date, total_views, total_engagements,
      total_content, total_creators, total_budget, avg_er, total_clicks,
      avg_cpc, creators (list of dicts with name, platform, views, likes,
      comments, link_clicks, er, avg_views, avg_ccv, peak_ccv)
    """
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # blank

    title     = campaign_data.get('campaign_title', 'Campaign Title')
    date_str  = campaign_data.get('campaign_date', 'Month Year')
    creators  = campaign_data.get('creators', [])

    # ── SLIDE 1: TITLE (dark with olive bottom wave) ──────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0x11, 0x11, 0x11)

    # Olive wave strip at bottom
    add_rect(s, 0, 6.2, 13.33, 1.3, OLIVE)
    add_rect(s, 0, 5.8, 13.33, 0.6, RGBColor(0x6B, 0x7D, 0x1A))

    # "cherrypick TALENT" bottom right
    add_text(s, "cherrypick", 9.5, 6.25, 3.5, 0.6,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(s, "TALENT", 9.5, 6.75, 3.5, 0.5,
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial")

    # Campaign title
    add_text(s, title, 0.7, 2.5, 8.5, 1.5,
             size=54, bold=True, color=WHITE, font_name="Arial Black")
    add_text(s, f"End of Campaign Report  |  {date_str}", 0.7, 4.1, 8.0, 0.6,
             size=20, bold=True, color=RGBColor(0xCC, 0xCC, 0xCC), font_name="Arial")

    # Confidentiality
    add_text(s, "This presentation and its contents are strictly confidential and intended solely for the recipient.",
             0.7, 6.95, 9.0, 0.35, size=8, color=RGBColor(0xAA, 0xAA, 0xAA), font_name="Arial")

    # ── SLIDE 2: CAMPAIGN SNAPSHOT (dark) ─────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    add_cpt_mark_white(s)

    add_text(s, "Campaign Snapshot", 0.5, 0.25, 9.0, 0.85,
             size=44, bold=True, color=WHITE, font_name="Arial Black")
    add_text(s, title, 0.5, 1.0, 9.0, 0.4,
             size=18, color=GREY, font_name="Arial")

    # 4 top stat cards + 4 bottom
    stats_top = [
        (fmt_num(campaign_data.get('total_views')),       "Total Views"),
        (fmt_num(campaign_data.get('total_content')),     "Content Pieces"),
        (fmt_num(campaign_data.get('total_creators')),    "Creators Booked"),
        (fmt_currency(campaign_data.get('total_budget')), "Total Budget"),
    ]
    stats_bot = [
        (fmt_num(campaign_data.get('total_engagements')), "Total Engagements"),
        (fmt_pct(campaign_data.get('avg_er')),            "Avg Engagement Rate"),
        (fmt_num(campaign_data.get('total_clicks')),      "Link Clicks"),
        (fmt_currency(campaign_data.get('avg_cpc')),      "Avg CPC"),
    ]
    card_w = 2.9
    for i, (val, lbl) in enumerate(stats_top):
        stat_card(s, 0.5 + i * 3.1, 1.6, card_w, 1.6, val, lbl, dark=False)
    for i, (val, lbl) in enumerate(stats_bot):
        stat_card(s, 0.5 + i * 3.1, 3.5, card_w, 1.6, val, lbl, dark=True)

    # Olive bottom bar
    add_rect(s, 0, 5.35, 13.33, 0.08, OLIVE)

    # ── SLIDE 3: OUR CREATORS (light) ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_cpt_mark(s)

    add_text(s, "Our Creators", 0.5, 0.2, 12.0, 1.0,
             size=44, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font_name="Arial Black")

    # Creator name grid — up to 15 names in 5 cols × 3 rows
    cols = 5
    names = [c.get('name', '—') for c in creators[:15]]
    card_w, card_h = 2.4, 1.5
    gap_x, gap_y = 0.2, 0.2
    start_x, start_y = 0.45, 1.4
    for idx, name in enumerate(names):
        col_i = idx % cols
        row_i = idx // cols
        cx = start_x + col_i * (card_w + gap_x)
        cy = start_y + row_i * (card_h + gap_y)
        add_rect(s, cx, cy, card_w, card_h, BLACK)
        add_rect(s, cx, cy + card_h - 0.38, card_w, 0.38, RGBColor(0x22, 0x22, 0x22))
        add_text(s, name, cx + 0.08, cy + card_h - 0.36, card_w - 0.16, 0.32,
                 size=11, bold=True, color=WHITE, font_name="Arial")

    # ── SLIDE 4: CAMPAIGN TIMELINE (light) ───────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_cpt_mark(s)

    add_text(s, "Campaign Timeline", 0.5, 0.2, 12.0, 0.85,
             size=44, bold=True, color=BLACK, font_name="Arial Black")
    add_rect(s, 0.5, 1.15, 12.33, 0.05, OLIVE)

    # Collect dated entries and sort chronologically
    from datetime import datetime as dt
    dated = []
    for c in creators:
        ld = c.get('live_date')
        if ld:
            try:
                # Try multiple date formats
                for fmt in ('%Y-%m-%d', '%m/%d/%Y', '%d/%m/%Y', '%Y-%m-%d %H:%M:%S'):
                    try:
                        parsed = dt.strptime(str(ld)[:10], fmt)
                        dated.append((parsed, c['name'], c.get('platform', '')))
                        break
                    except: continue
            except: pass

    dated.sort(key=lambda x: x[0])

    if dated:
        # Group by date
        from collections import defaultdict
        by_date = defaultdict(list)
        for d, name, plat in dated:
            by_date[d].append((name, plat))

        sorted_dates = sorted(by_date.keys())
        n = len(sorted_dates)

        # Draw horizontal timeline spine
        spine_y = 3.4
        add_rect(s, 0.7, spine_y - 0.02, 11.93, 0.06, OLIVE)

        # Distribute date stops evenly
        usable_w = 11.5
        step = usable_w / max(n - 1, 1) if n > 1 else 0
        start_x = 0.9

        for i, date_key in enumerate(sorted_dates):
            names_here = by_date[date_key]
            cx = start_x + i * step
            above = (i % 2 == 0)  # alternate above/below

            # Dot on spine
            add_rect(s, cx - 0.08, spine_y - 0.1, 0.18, 0.18,
                     OLIVE if above else BLACK)

            # Date label
            date_label = date_key.strftime('%b %d')
            if above:
                add_text(s, date_label, cx - 0.55, spine_y - 0.55, 1.2, 0.3,
                         size=9, bold=True, color=OLIVE, align=PP_ALIGN.CENTER, font_name="Arial")
            else:
                add_text(s, date_label, cx - 0.55, spine_y + 0.2, 1.2, 0.3,
                         size=9, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font_name="Arial")

            # Creator names (up to 3 per date, then "+X more")
            display = names_here[:3]
            extra = len(names_here) - 3
            for j, (nm, _) in enumerate(display):
                if above:
                    ny = spine_y - 1.0 - j * 0.32
                else:
                    ny = spine_y + 0.55 + j * 0.32
                add_text(s, nm, cx - 0.7, ny, 1.5, 0.28,
                         size=8, color=DARK_GREY, align=PP_ALIGN.CENTER, font_name="Arial")
            if extra > 0:
                ey = spine_y - 1.0 - 3 * 0.32 if above else spine_y + 0.55 + 3 * 0.32
                add_text(s, f"+{extra} more", cx - 0.7, ey, 1.5, 0.28,
                         size=8, color=GREY, align=PP_ALIGN.CENTER, font_name="Arial")

        # Summary stats below
        total_dated = len(dated)
        if sorted_dates:
            span_days = (sorted_dates[-1] - sorted_dates[0]).days
            first_date = sorted_dates[0].strftime('%b %d, %Y')
            last_date = sorted_dates[-1].strftime('%b %d, %Y')
            summary = f"{total_dated} pieces of content went live across {span_days + 1} days  ·  {first_date} → {last_date}"
        else:
            summary = f"{total_dated} pieces of content"

        add_rect(s, 0.5, 6.7, 12.33, 0.55, CREAM)
        add_text(s, summary, 0.7, 6.75, 11.9, 0.42,
                 size=11, color=DARK_GREY, align=PP_ALIGN.CENTER, font_name="Arial")
    else:
        # No dates — placeholder
        add_rect(s, 0.5, 2.5, 12.33, 2.0, CREAM)
        add_text(s, "Add live dates to the LIVE DATE column in your template to auto-generate the timeline.",
                 1.0, 3.1, 11.33, 0.6,
                 size=14, color=GREY, align=PP_ALIGN.CENTER, font_name="Arial")

    # ── SLIDE 5: PERFORMANCE SUMMARY (dark) ──────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    add_cpt_mark_white(s)

    add_text(s, "Performance Summary", 0.5, 0.2, 12.0, 0.9,
             size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")

    # 4 big stat boxes
    perf_stats = [
        (fmt_num(campaign_data.get('total_views')),       "Total Views Delivered"),
        (fmt_num(campaign_data.get('total_engagements')), "Total Engagements"),
        (fmt_pct(campaign_data.get('avg_er')),            "Campaign Avg ER"),
        (fmt_num(campaign_data.get('total_clicks')),      "Total Link Clicks"),
    ]
    bw = 2.9
    for i, (val, lbl) in enumerate(perf_stats):
        bx = 0.5 + i * 3.1
        add_rect(s, bx, 1.4, bw, 2.6, WHITE)
        add_text(s, val, bx + 0.1, 1.6, bw - 0.2, 1.4,
                 size=40, bold=True, color=BLACK, align=PP_ALIGN.CENTER, font_name="Arial Black")
        add_text(s, lbl, bx + 0.1, 3.1, bw - 0.2, 0.6,
                 size=13, color=DARK_GREY, align=PP_ALIGN.CENTER, font_name="Arial")

    # Olive accent line
    add_rect(s, 0, 4.25, 13.33, 0.08, OLIVE)

    # ── SLIDE 6: CREATOR BREAKDOWN TABLE (light) ──────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = CREAM
    add_cpt_mark(s)

    add_text(s, "Creator Breakdown", 0.5, 0.2, 12.0, 0.85,
             size=44, bold=True, color=BLACK, font_name="Arial Black")
    add_rect(s, 0.5, 1.15, 12.33, 0.05, OLIVE)

    # Table headers
    headers = ["Creator", "Platform", "Views", "Likes", "Comments", "Link Clicks", "ER %"]
    col_widths = [2.8, 1.5, 1.4, 1.2, 1.4, 1.5, 1.0]
    hx = 0.5
    for hdr, cw in zip(headers, col_widths):
        add_rect(s, hx, 1.25, cw - 0.05, 0.4, BLACK)
        add_text(s, hdr, hx + 0.08, 1.28, cw - 0.15, 0.35,
                 size=10, bold=True, color=WHITE, font_name="Arial")
        hx += cw

    # Table rows — up to 12 creators
    row_h = 0.42
    for r_idx, creator in enumerate(creators[:12]):
        ry = 1.7 + r_idx * row_h
        bg = WHITE if r_idx % 2 == 0 else CREAM
        add_rect(s, 0.5, ry, 12.33, row_h - 0.03, bg)

        vals = [
            creator.get('name', '—'),
            creator.get('platform', '—'),
            fmt_num(creator.get('views')),
            fmt_num(creator.get('likes')),
            fmt_num(creator.get('comments')),
            fmt_num(creator.get('link_clicks')),
            fmt_pct(creator.get('er')),
        ]
        rx = 0.5
        for val, cw in zip(vals, col_widths):
            add_text(s, val, rx + 0.08, ry + 0.06, cw - 0.15, row_h - 0.1,
                     size=10, color=BLACK, font_name="Arial")
            rx += cw

    # ── SLIDE 7: TOP PERFORMERS (dark) ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    add_cpt_mark_white(s)

    add_text(s, "Top Performers", 0.5, 0.2, 12.0, 0.85,
             size=44, bold=True, color=WHITE, font_name="Arial Black")
    add_text(s, "Highest performing creators by views this campaign", 0.5, 1.0, 12.0, 0.4,
             size=16, color=GREY, font_name="Arial")

    # Sort by views and show top 5
    sorted_creators = sorted(
        [c for c in creators if c.get('views') and str(c.get('views')) not in ('', '-', 'None')],
        key=lambda x: float(str(x.get('views', 0)).replace(',', '') or 0),
        reverse=True
    )[:5]

    card_w = 2.3
    for i, c in enumerate(sorted_creators):
        cx = 0.5 + i * (card_w + 0.2)
        add_rect(s, cx, 1.6, card_w, 3.5, RGBColor(0x22, 0x22, 0x22))
        add_rect(s, cx, 1.6, card_w, 0.08, OLIVE)
        rank_label = ["🥇", "🥈", "🥉", "4th", "5th"][i]
        add_text(s, rank_label, cx + 0.1, 1.65, card_w - 0.2, 0.45,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial")
        add_text(s, c.get('name', '—'), cx + 0.1, 2.15, card_w - 0.2, 0.45,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial")
        add_text(s, c.get('platform', ''), cx + 0.1, 2.55, card_w - 0.2, 0.3,
                 size=10, color=GREY, align=PP_ALIGN.CENTER, font_name="Arial")
        add_rect(s, cx + 0.15, 2.95, card_w - 0.3, 0.03, OLIVE)
        metrics = [
            ("Views",    fmt_num(c.get('views'))),
            ("Likes",    fmt_num(c.get('likes'))),
            ("Comments", fmt_num(c.get('comments'))),
            ("Clicks",   fmt_num(c.get('link_clicks'))),
            ("ER",       fmt_pct(c.get('er'))),
        ]
        for m_idx, (mlbl, mval) in enumerate(metrics):
            my = 3.1 + m_idx * 0.38
            add_text(s, mlbl, cx + 0.12, my, 1.0, 0.32,
                     size=9, color=GREY, font_name="Arial")
            add_text(s, mval, cx + 0.12, my, card_w - 0.25, 0.32,
                     size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, font_name="Arial")

    # ── SLIDE 8: LINK PERFORMANCE (light) ─────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    add_cpt_mark(s)

    add_text(s, "Link Performance", 0.5, 0.2, 10.0, 0.85,
             size=44, bold=True, color=BLACK, font_name="Arial Black")
    add_rect(s, 0.5, 1.15, 12.33, 0.05, OLIVE)

    # Summary stats row
    link_stats = [
        ("Total Clicks",   fmt_num(campaign_data.get('total_clicks'))),
        ("Total Budget",   fmt_currency(campaign_data.get('total_budget'))),
        ("Avg CPC",        fmt_currency(campaign_data.get('avg_cpc'))),
    ]
    for i, (lbl, val) in enumerate(link_stats):
        bx = 0.5 + i * 4.2
        add_rect(s, bx, 1.35, 3.8, 1.4, BLACK)
        add_text(s, val, bx + 0.15, 1.45, 3.5, 0.8,
                 size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
        add_text(s, lbl, bx + 0.15, 2.2, 3.5, 0.4,
                 size=12, color=GREY, align=PP_ALIGN.CENTER, font_name="Arial")

    # Per-creator link table
    headers2 = ["Creator", "Platform", "Link Clicks", "Budget", "CPC"]
    col_widths2 = [3.5, 2.0, 2.2, 2.2, 2.0]
    hx = 0.5
    for hdr, cw in zip(headers2, col_widths2):
        add_rect(s, hx, 3.0, cw - 0.05, 0.38, BLACK)
        add_text(s, hdr, hx + 0.08, 3.03, cw - 0.15, 0.32,
                 size=10, bold=True, color=WHITE, font_name="Arial")
        hx += cw

    link_creators = [c for c in creators if c.get('link_clicks') and
                     str(c.get('link_clicks')) not in ('', '-', 'None', '0')]
    for r_idx, creator in enumerate(link_creators[:10]):
        ry = 3.43 + r_idx * 0.38
        bg = WHITE if r_idx % 2 == 0 else CREAM
        add_rect(s, 0.5, ry, 12.33, 0.35, bg)
        vals2 = [
            creator.get('name', '—'),
            creator.get('platform', '—'),
            fmt_num(creator.get('link_clicks')),
            fmt_currency(creator.get('budget')),
            fmt_currency(creator.get('cpc')),
        ]
        rx = 0.5
        for val, cw in zip(vals2, col_widths2):
            add_text(s, val, rx + 0.08, ry + 0.05, cw - 0.15, 0.28,
                     size=10, color=BLACK, font_name="Arial")
            rx += cw

    # ── SLIDE 9: CAMPAIGN LEARNINGS (dark) ────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BLACK
    add_cpt_mark_white(s)

    add_text(s, "Campaign Learnings", 0.5, 0.2, 12.0, 0.85,
             size=44, bold=True, color=WHITE, font_name="Arial Black")
    add_rect(s, 0.5, 1.15, 12.33, 0.05, OLIVE)

    # 3 learning placeholders
    for i in range(3):
        ly = 1.4 + i * 1.8
        add_rect(s, 0.5, ly, 12.33, 1.6, RGBColor(0x1A, 0x1A, 0x1A))
        add_rect(s, 0.5, ly, 0.12, 1.6, OLIVE)
        add_text(s, f"Learning {i+1}", 0.8, ly + 0.15, 11.5, 0.4,
                 size=16, bold=True, color=WHITE, font_name="Arial")
        add_text(s, "Add your campaign learning here — what worked, what didn't, and recommendations for next time.",
                 0.8, ly + 0.55, 11.5, 0.85,
                 size=13, color=GREY, font_name="Arial")

    # ── SLIDE 10: THANK YOU (dark with wave) ───────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = RGBColor(0x33, 0x3D, 0x2A)

    add_rect(s, 0, 5.2, 13.33, 2.3, RGBColor(0x2D, 0x38, 0x06))
    add_rect(s, 0, 4.8, 13.33, 0.6, RGBColor(0x4A, 0x5A, 0x0A))

    add_text(s, "Thank You!", 0, 2.0, 13.33, 1.5,
             size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(s, "lead@cherrypicktalent.com", 0, 3.5, 13.33, 0.5,
             size=18, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER, font_name="Arial")
    add_text(s, "cherrypick  TALENT", 0, 5.8, 13.33, 0.7,
             size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Georgia")

    prs.save(output_path)
    return output_path


def build_campaign_data_from_xlsx(xlsx_path, sheet_name=None):
    """Read the filled campaign report xlsx and extract summary data for PPTX."""
    from openpyxl import load_workbook
    wb = load_workbook(xlsx_path, data_only=True)
    ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

    # Find header row
    header_row = None
    for i, row in enumerate(ws.iter_rows(values_only=True), 1):
        if row[0] and str(row[0]).upper() == 'TALENT':
            header_row = i
            break
    if not header_row:
        return {}

    # Fixed column positions
    # Auto-detect columns — handles old and new template versions
    keywords = {
        'talent':           ['TALENT'],
        'live_date':        ['LIVE DATE'],
        'content_url':      ['LINK TO CONTENT', 'CONTENT URL'],
        'rebrandly':        ['REBRANDLY'],
        'campaign':         ['CAMPAIGN', 'GAME'],
        'format':           ['FORMAT'],
        'platform':         ['PLATFORM'],
        'followers':        ['FOLLOWERS'],
        'views':            ['VIEWS'],
        'avg_views':        ['AVG VIEWS', '7 DAY AVG', 'AVG. VIEWS'],
        'avg_ccv':          ['AVG. CCV', 'AVG CCV'],
        'peak_ccv':         ['PEAK CCV'],
        'hours':            ['HOURS WATCHED'],
        'likes':            ['LIKES'],
        'comments':         ['COMMENTS'],
        'shares':           ['SHARES'],
        'saves':            ['SAVES'],
        'total_engagement': ['TOTAL ENGAGEMENT'],
        'engagement_rate':  ['ENGAGEMENT RATE'],
        'link_clicks':      ['LINK CLICKS'],
        'pct_goal':         ['% TO'],
        'budget':           ['BUDGET'],
        'cpc':              ['CPC'],
    }
    COL = {}
    for cell in ws[header_row]:
        if not cell.value: continue
        h = str(cell.value).strip().upper().replace('\n', ' ')
        for key, terms in keywords.items():
            if key not in COL and any(t in h for t in terms):
                COL[key] = cell.column - 1
                break

    def safe_float(v):
        if v is None or str(v).strip() in ('', '-', 'None'): return None
        try: return float(str(v).replace(',','').replace('$','').replace('%',''))
        except: return None

    creators = []
    total_views = total_engagements = total_clicks = total_budget = 0
    er_values = []
    cpc_values = []
    campaign_name = None

    for row in ws.iter_rows(min_row=header_row+1, values_only=True):
        if not row[0] or str(row[0]).strip() in ('', 'TALENT'):
            continue
        name     = str(row[COL['talent']] or '').strip()
        platform = str(row[COL['platform']] or '').strip()
        _camp = COL.get('campaign'); campaign_name = campaign_name or str((row[_camp] if _camp is not None else None) or '').strip()

        def gcol(key):
            idx = COL.get(key)
            return row[idx] if idx is not None else None

        live_date = gcol('live_date')
        views    = safe_float(gcol('views'))
        likes    = safe_float(gcol('likes'))
        comments = safe_float(gcol('comments'))
        clicks   = safe_float(gcol('link_clicks'))
        eng      = safe_float(gcol('total_engagement'))
        er_raw   = safe_float(gcol('engagement_rate'))
        budget   = safe_float(gcol('budget'))
        cpc      = safe_float(gcol('cpc'))
        avg_ccv  = safe_float(gcol('avg_ccv'))
        peak_ccv = safe_float(gcol('peak_ccv'))

        # Normalise live_date to a string
        if live_date is None or str(live_date).strip() in ('', 'None', '-'):
            live_date_str = None
        elif hasattr(live_date, 'strftime'):
            live_date_str = live_date.strftime('%Y-%m-%d')
        else:
            live_date_str = str(live_date).strip()[:10]

        creators.append({
            'name': name, 'platform': platform,
            'live_date': live_date_str,
            'views': int(views) if views else None,
            'likes': int(likes) if likes else None,
            'comments': int(comments) if comments else None,
            'link_clicks': int(clicks) if clicks else None,
            'er': round(er_raw, 2) if er_raw else None,
            'budget': budget, 'cpc': cpc,
            'avg_ccv': avg_ccv, 'peak_ccv': peak_ccv,
        })

        if views:   total_views        += views
        if eng:     total_engagements  += eng
        if clicks:  total_clicks       += clicks
        if budget:  total_budget       += budget
        if er_raw:  er_values.append(er_raw)
        if cpc:     cpc_values.append(cpc)

    avg_er  = round(sum(er_values) / len(er_values), 2) if er_values else None
    avg_cpc = round(sum(cpc_values) / len(cpc_values), 2) if cpc_values else None

    from datetime import datetime
    return {
        'campaign_title': campaign_name or 'Campaign Report',
        'campaign_date': datetime.now().strftime('%B %Y'),
        'total_views': int(total_views),
        'total_engagements': int(total_engagements),
        'total_content': len(creators),
        'total_creators': len(set(c['name'] for c in creators)),
        'total_budget': round(total_budget, 2),
        'avg_er': avg_er,
        'total_clicks': int(total_clicks),
        'avg_cpc': avg_cpc,
        'creators': creators,
    }
